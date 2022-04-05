# coding: UTF-8

from pptx import Presentation
from pdf2image import convert_from_path, convert_from_bytes
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
import os
import datetime
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

dt_name = datetime.datetime.now()
date = '%d-%d-%d' % (dt_name.year,dt_name.month,dt_name.day)

for filename in os.listdir('source_files/'):
    if os.path.splitext(filename)[1] == '.pdf':
        print("Creating %s" % filename)
        prs = Presentation()

        pages = convert_from_path('source_files/' + filename, 500)
        for index, page in enumerate(pages):
            #Save as 'jpg' in jpgs dir
            jpg_file = "jpgs/%s-(%d).jpg" % (filename,index)
            page.save(jpg_file, 'JPEG')

            #Get width/height of image
            image = Image.open(jpg_file)
            height = image.height
            width = image.width
            #Rotate 270 degrees if horizontal
            if height > width:
                adjusted = image.rotate(270, expand=True)
                adjusted.save(jpg_file)

            #Setup slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            left = top = 0
            slide.shapes.add_picture(jpg_file, left,top,height = prs.slide_height)

        prs.save('result/%s.pptx' % os.path.splitext(filename)[0])

    else:
        print("Skipping %s because it\'s not a pdf" % filename)

print("Saved to result directory")
